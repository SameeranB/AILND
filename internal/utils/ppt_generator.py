import re
from typing import List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.text.text import _Run
import streamlit as st
import io

def format_text_run(run: _Run, text: str):
    """Apply formatting to a text run based on markdown syntax."""
    # Remove the formatting markers
    if text.startswith('**') and text.endswith('**'):
        text = text[2:-2]
        run.font.bold = True
    elif text.startswith('*') and text.endswith('*'):
        text = text[1:-1]
        run.font.italic = True
    elif text.startswith('_') and text.endswith('_'):
        text = text[1:-1]
        run.font.italic = True
    elif text.startswith('####'):
        text = text[4:].strip()
        run.font.size = Pt(14)
        run.font.bold = True
    run.text = text

def format_paragraph_text(paragraph, text: str):
    """Parse and format text with markdown syntax."""
    # Set default paragraph level to 0 (no bullets)
    paragraph.level = 0
    
    # Split text into formatting segments
    segments = []
    current_pos = 0
    
    # Find all formatting patterns
    patterns = [
        (r'\*\*.*?\*\*', 'bold'),     # Bold
        (r'_.*?_', 'italic'),         # Italic (underscore style)
        (r'####.*?($|\n)', 'h4')      # Level 4 heading
    ]
    
    # Find all matches for all patterns
    matches = []
    for pattern, fmt_type in patterns:
        for match in re.finditer(pattern, text):
            matches.append((match.start(), match.end(), match.group(), fmt_type))
    
    # Sort matches by start position
    matches.sort(key=lambda x: x[0])
    
    # Process matches and non-formatted text
    current_pos = 0
    for start, end, content, fmt_type in matches:
        # Add non-formatted text before the match
        if start > current_pos:
            segments.append((text[current_pos:start], 'normal'))
        segments.append((content, fmt_type))
        current_pos = end
    
    # Add remaining text
    if current_pos < len(text):
        segments.append((text[current_pos:], 'normal'))
    
    # Apply formatting to each segment
    for text_segment, fmt_type in segments:
        run = paragraph.add_run()
        if fmt_type == 'bold':
            run.font.bold = True
            run.text = text_segment[2:-2]  # Remove ** markers
        elif fmt_type == 'italic':
            run.font.italic = True
            run.text = text_segment[1:-1]  # Remove _ markers
        elif fmt_type == 'h4':
            run.font.bold = True
            run.font.size = Pt(14)
            run.text = text_segment[4:].strip()  # Remove #### marker
        else:
            run.text = text_segment

def parse_markdown_content(content: str) -> List[Tuple[str, str]]:
    """
    Parse markdown content into slides, where each slide has a title and content.
    Returns a list of tuples (title, content).
    """
    # First split content by horizontal rule separator
    slide_chunks = re.split(r'\n\s*---\s*\n', content)
    
    slides = []
    
    for chunk in slide_chunks:
        # Split chunk into lines
        lines = chunk.strip().split('\n')
        current_title = None
        current_content = []
        
        for line in lines:
            # Skip horizontal rule lines
            if re.match(r'^\s*---\s*$', line):
                continue
                
            # Check if line is a heading (h1, h2, h3)
            heading_match = re.match(r'^(#{1,3})\s+(.+)$', line)
            if heading_match:
                # If we have a previous slide, save it
                if current_title:
                    slides.append((current_title, '\n'.join(current_content).strip()))
                # Start new slide
                current_title = heading_match.group(2)
                current_content = []
            else:
                # Add line to current content if we have a title
                if current_title and line.strip():
                    current_content.append(line)
        
        # Add the last slide from this chunk if exists
        if current_title and current_content:
            slides.append((current_title, '\n'.join(current_content).strip()))
    
    return slides

def create_presentation(slides: List[Tuple[str, str]]) -> Presentation:
    """
    Create a PowerPoint presentation from the parsed slides.
    """
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content
    
    # Create title slide
    if slides:
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        title.text = slides[0][0]
        
        # Add subtitle if there's content in first slide
        if slides[0][1]:
            subtitle = title_slide.placeholders[1]
            # Format the subtitle text
            paragraph = subtitle.text_frame.paragraphs[0]
            format_paragraph_text(paragraph, slides[0][1])
    
    # Create content slides
    for title, content in slides[1:]:
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set content with formatting
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.clear()  # Clear existing paragraphs
        
        # Split content into paragraphs
        paragraphs = content.split('\n')
        
        for i, para_text in enumerate(paragraphs):
            # Skip empty lines
            if not para_text.strip():
                continue
                
            # Create new paragraph
            if i == 0:
                paragraph = tf.paragraphs[0]
            else:
                paragraph = tf.add_paragraph()
            
            # Format the paragraph text (this will also set the correct level)
            format_paragraph_text(paragraph, para_text)
            
            # Set paragraph properties
            paragraph.alignment = PP_ALIGN.LEFT
    
    return prs

def generate_ppt_from_markdown(content: str) -> bytes:
    """
    Generate a PowerPoint presentation from markdown content.
    Returns the presentation as bytes.
    """
    # Parse markdown into slides
    slides = parse_markdown_content(content)
    
    if not slides:
        raise ValueError("No valid slides found in the content")
    
    # Create presentation
    prs = create_presentation(slides)
    
    # Save to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes.getvalue() 